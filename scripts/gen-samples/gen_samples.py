"""Generates the minimal sample files used by the integration test suite.

One file per supported format + one rejected format. Each sample carries a
distinctive marker string the tests look for in the converted markdown.

Run with:
    uv run --with python-docx --with python-pptx --with openpyxl --with xlwt \
        python gen_samples.py <output-dir>
"""

import io
import struct
import sys
import zipfile
from pathlib import Path

OUT = Path(sys.argv[1] if len(sys.argv) > 1 else "samples")
OUT.mkdir(parents=True, exist_ok=True)

MARKER = "mdNotch integration sample"


def write(name: str, data: bytes) -> None:
    (OUT / name).write_bytes(data)
    print(f"wrote {name} ({len(data)} bytes)")


# --- PDF (hand-built, no deps) -------------------------------------------
def gen_pdf() -> None:
    text = f"{MARKER} PDF"
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R"
        b" /Resources << /Font << /F1 5 0 R >> >> >>",
    ]
    stream = b"BT /F1 24 Tf 72 700 Td (%s) Tj ET" % text.encode()
    objs.append(b"<< /Length %d >>\nstream\n%s\nendstream" % (len(stream), stream))
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (i, o)
    xref = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objs) + 1,
        xref,
    )
    write("sample.pdf", bytes(out))


# --- Plain-text formats ---------------------------------------------------
def gen_text_formats() -> None:
    write(
        "sample.html",
        f"<html><body><h1>{MARKER} HTML</h1><p>paragraph</p></body></html>".encode(),
    )
    write("sample.csv", f"name,value\n{MARKER} CSV,42\n".encode())
    write("sample.json", ('{"marker": "%s JSON", "value": 42}' % MARKER).encode())
    write(
        "sample.xml",
        f'<?xml version="1.0"?><root><item>{MARKER} XML</item></root>'.encode(),
    )


# --- Office formats (real libraries so the readers accept them) ----------
def gen_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph(f"{MARKER} DOCX")
    buf = io.BytesIO()
    doc.save(buf)
    write("sample.docx", buf.getvalue())


def gen_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"{MARKER} PPTX"
    buf = io.BytesIO()
    prs.save(buf)
    write("sample.pptx", buf.getvalue())


def gen_xlsx() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.active["A1"] = f"{MARKER} XLSX"
    buf = io.BytesIO()
    wb.save(buf)
    write("sample.xlsx", buf.getvalue())


def gen_xls() -> None:
    import xlwt

    wb = xlwt.Workbook()
    sheet = wb.add_sheet("Sheet1")
    sheet.write(0, 0, f"{MARKER} XLS")
    buf = io.BytesIO()
    wb.save(buf)
    write("sample.xls", buf.getvalue())


# --- EPUB (hand-built zip) -----------------------------------------------
def gen_epub() -> None:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("mimetype", "application/epub+zip", compress_type=zipfile.ZIP_STORED)
        z.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?>'
            '<container version="1.0" xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="OEBPS/content.opf"'
            ' media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        z.writestr(
            "OEBPS/content.opf",
            '<?xml version="1.0"?>'
            '<package xmlns="http://www.idpf.org/2007/opf" version="2.0"'
            ' unique-identifier="id">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            "<dc:title>mdNotch sample epub</dc:title>"
            '<dc:identifier id="id">mdnotch-sample</dc:identifier>'
            "<dc:language>en</dc:language></metadata>"
            '<manifest><item id="ch1" href="chapter1.xhtml"'
            ' media-type="application/xhtml+xml"/></manifest>'
            '<spine><itemref idref="ch1"/></spine></package>',
        )
        z.writestr(
            "OEBPS/chapter1.xhtml",
            '<?xml version="1.0"?>'
            '<html xmlns="http://www.w3.org/1999/xhtml"><head>'
            "<title>Chapter 1</title></head>"
            f"<body><h1>{MARKER} EPUB</h1></body></html>",
        )
    write("sample.epub", buf.getvalue())


# --- ZIP: two inner files, checks recursive concatenation ----------------
def gen_zip() -> None:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("inner-one.csv", f"col\n{MARKER} ZIP-inner-one\n")
        z.writestr(
            "inner-two.html",
            f"<html><body><p>{MARKER} ZIP-inner-two</p></body></html>",
        )
    write("sample.zip", buf.getvalue())


# --- Rejected format: 1x1 PNG --------------------------------------------
def gen_png() -> None:
    import binascii

    def chunk(typ: bytes, data: bytes) -> bytes:
        c = struct.pack(">I", len(data)) + typ + data
        return c + struct.pack(">I", binascii.crc32(typ + data) & 0xFFFFFFFF)

    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 0, 0, 0, 0)
    idat = b"\x78\x9c\x62\x60\x00\x00\x00\x02\x00\x01"  # zlib: one gray pixel
    png = (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", idat)
        + chunk(b"IEND", b"")
    )
    write("sample.png", png)


gen_pdf()
gen_text_formats()
gen_docx()
gen_pptx()
gen_xlsx()
gen_xls()
gen_epub()
gen_zip()
gen_png()
print("done")
